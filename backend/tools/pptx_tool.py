"""
PowerPoint Generator Tool.

Generates a structured .pptx report from technical research findings.
"""
from __future__ import annotations

import os
import uuid
import logging
from typing import List
from pydantic import BaseModel, Field
from copilot import define_tool

logger = logging.getLogger(__name__)

OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), "generated_reports")


class Slide(BaseModel):
    title: str = Field(description="Slide title")
    bullets: List[str] = Field(description="Bullet points for this slide")


class GeneratePowerpointParams(BaseModel):
    report_title: str = Field(description="Title for the PowerPoint report")
    overview: str = Field(description="Executive overview / summary of the issue")
    root_cause: str = Field(description="Root cause analysis")
    technical_details: str = Field(description="Technical deep-dive details")
    recommendations: str = Field(description="Recommended remediation steps")
    references: List[str] = Field(
        default_factory=list,
        description="List of reference URLs or documentation links",
    )


@define_tool(
    description=(
        "Generate a PowerPoint (.pptx) report summarizing technical research findings. "
        "Creates a structured 5-slide deck: Overview, Root Cause, Technical Deep Dive, "
        "Recommendations, and References. Returns a download path for the generated file."
    )
)
async def generate_powerpoint_tool(params: GeneratePowerpointParams) -> str:
    """Generate a .pptx report from structured research findings."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return (
            "[PowerPoint] python-pptx is not installed. "
            "Run: pip install python-pptx"
        )

    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ----- Slide layout helpers -----
    TITLE_SLIDE_LAYOUT = prs.slide_layouts[0]
    CONTENT_LAYOUT = prs.slide_layouts[1]

    DARK_BLUE = RGBColor(0x00, 0x3A, 0x70)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    def _add_title_slide(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(TITLE_SLIDE_LAYOUT)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = WHITE
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
        slide.placeholders[1].text = subtitle

    def _add_content_slide(title: str, body_text: str) -> None:
        slide = prs.slides.add_slide(CONTENT_LAYOUT)
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)

        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(body_text.strip().split("\n")):
            if i == 0:
                tf.paragraphs[0].text = line
            else:
                p = tf.add_paragraph()
                p.text = line
            tf.paragraphs[i].font.size = Pt(16)

    def _add_references_slide(refs: List[str]) -> None:
        slide = prs.slides.add_slide(CONTENT_LAYOUT)
        slide.shapes.title.text = "References"
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = DARK_BLUE
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)

        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, ref in enumerate(refs or ["No references provided."]):
            if i == 0:
                tf.paragraphs[0].text = ref
            else:
                p = tf.add_paragraph()
                p.text = ref
            tf.paragraphs[i].font.size = Pt(14)

    # ----- Build slides -----
    _add_title_slide(params.report_title, "Microsoft Support Agent | Technical Research Report")
    _add_content_slide("📋 Overview", params.overview)
    _add_content_slide("🔍 Root Cause Analysis", params.root_cause)
    _add_content_slide("⚙️ Technical Deep Dive", params.technical_details)
    _add_content_slide("✅ Recommended Actions", params.recommendations)
    _add_references_slide(params.references)

    filename = f"report_{uuid.uuid4().hex[:8]}.pptx"
    filepath = os.path.join(OUTPUT_DIR, filename)
    prs.save(filepath)

    logger.info("PowerPoint report saved to %s", filepath)
    return (
        f"PowerPoint report generated successfully.\n"
        f"File: {filepath}\n"
        f"Download via: GET /reports/{filename}"
    )
